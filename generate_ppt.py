from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_LINE_DASH_STYLE, MSO_ARROWHEAD_STYLE
from pptx.dml.color import RGBColor

TITLE_FONT_SIZE = Pt(24)
BOX_FONT_SIZE = Pt(16)
LABEL_FONT_SIZE = Pt(12)

LAVENDER = RGBColor(235, 232, 255)
BORDER = RGBColor(140, 114, 219)
PALE_YELLOW = RGBColor(255, 252, 220)
GRAY = RGBColor(90, 90, 90)

def add_title(slide, text):
    title = slide.shapes.title
    if title is None:
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = TITLE_FONT_SIZE
    run.font.bold = True

def add_panel(slide, x, y, w, h, caption=None):
    panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    panel.fill.solid()
    panel.fill.fore_color.rgb = PALE_YELLOW
    panel.line.color.rgb = BORDER
    panel.line.width = Pt(1.25)
    panel.shadow.inherit = False
    if caption:
        cap = slide.shapes.add_textbox(x + Inches(0.25), y - Inches(0.35), Inches(5), Inches(0.3))
        tf = cap.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = caption
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = GRAY
    return panel

def add_box(slide, x, y, w, h, text):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = LAVENDER
    shp.line.color.rgb = BORDER
    shp.line.width = Pt(1.25)
    tf = shp.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = BOX_FONT_SIZE
    r.font.color.rgb = RGBColor(20, 20, 20)
    r.font.bold = False
    return shp

def add_diamond(slide, x, y, w, h, text):
    shp = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = LAVENDER
    shp.line.color.rgb = BORDER
    shp.line.width = Pt(1.25)
    tf = shp.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = BOX_FONT_SIZE
    r.font.color.rgb = RGBColor(20, 20, 20)
    return shp

def mid_right(shape):
    return (shape.left + shape.width, shape.top + shape.height // 2)

def mid_left(shape):
    return (shape.left, shape.top + shape.height // 2)

def top_center(shape):
    return (shape.left + shape.width // 2, shape.top)

def bottom_center(shape):
    return (shape.left + shape.width // 2, shape.top + shape.height)

def add_arrow(slide, x1, y1, x2, y2, width_pt=2.0, dashed=False):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = GRAY
    conn.line.width = Pt(width_pt)
    conn.line.end_arrowhead.style = MSO_ARROWHEAD_STYLE.TRIANGLE
    if dashed:
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return conn

def add_label(slide, x, y, w, h, text):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = LABEL_FONT_SIZE
    r.font.color.rgb = GRAY
    return tb

def connect_lr(slide, left_shape, right_shape, dashed=False):
    x1, y1 = mid_right(left_shape)
    x2, y2 = mid_left(right_shape)
    return add_arrow(slide, x1, y1, x2, y2, dashed=dashed)

def elbow_feedback(slide, from_shape, to_shape, via_y):
    x_up = from_shape.left + from_shape.width // 2
    add_arrow(slide, x_up, from_shape.top, x_up, via_y)
    x_to = to_shape.left + to_shape.width // 2
    add_arrow(slide, x_up, via_y, x_to, via_y)
    add_arrow(slide, x_to, via_y, x_to, to_shape.top)

def slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    add_title(slide, "Pure Data-Driven Approach: Macroscale (Horizontal)")
    add_panel(slide, Inches(0.25), Inches(0.9), Inches(9.5), Inches(5.8))
    y = Inches(3.1)
    box_w, box_h, gap = Inches(2.0), Inches(1.0), Inches(0.35)
    x = Inches(0.4)

    s_exp = add_box(slide, x, y, box_w, box_h, "Macroscale Experiments"); x += box_w + gap
    s_data = add_box(slide, x, y, box_w, box_h, "Macro Experimental Data"); x += box_w + gap
    s_core = add_box(slide, x, y, box_w, box_h, "AI Learning Core\nSINDy / Bayesian\nCalibration"); x += box_w + gap
    s_valm = add_box(slide, x, y, box_w, box_h, "Validated Predictive Model"); x += box_w + gap
    s_oed  = add_diamond(slide, x, y - Inches(0.1), Inches(1.6), Inches(1.6), "Design of New\nExperiments")

    connect_lr(slide, s_exp, s_data)
    connect_lr(slide, s_data, s_core)
    connect_lr(slide, s_core, s_valm)
    x1, y1 = mid_right(s_valm)
    x2, y2 = s_oed.left, s_oed.top + s_oed.height // 2
    add_arrow(slide, x1, y1, x2, y2)

    elbow_feedback(slide, s_oed, s_exp, via_y=Inches(1.4))

def slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    add_title(slide, "Hybrid Model: Macro (Top) → Meso (Bottom), Both Horizontal")

    macro_panel = add_panel(slide, Inches(0.3), Inches(1.0), Inches(6.8), Inches(3.0), "Macroscale Learning Loop")
    meso_panel  = add_panel(slide, Inches(0.8), Inches(4.0), Inches(8.6), Inches(3.0), "Mesoscale Learning Loop")

    yM = Inches(2.1)
    w, h, g = Inches(1.7), Inches(0.9), Inches(0.3)
    xM = macro_panel.left + Inches(0.3)

    m_exp = add_box(slide, xM, yM, w, h, "Macro Experiments"); xM += w + g
    m_data = add_box(slide, xM, yM, w, h, "Macro Data"); xM += w + g
    m_core = add_box(slide, xM, yM, w, h, "Macro AI Core"); xM += w + g
    m_val  = add_box(slide, xM, yM, w, h, "Validated Macro Model"); xM += w + g
    m_oed  = add_diamond(slide, xM, yM - Inches(0.05), Inches(1.2), Inches(1.2), "Macro OED")

    connect_lr(slide, m_exp, m_data)
    connect_lr(slide, m_data, m_core)
    connect_lr(slide, m_core, m_val)
    connect_lr(slide, m_val, m_oed)
    elbow_feedback(slide, m_oed, m_exp, via_y=macro_panel.top + Inches(1.15))

    yS = Inches(5.2)
    w2, h2, g2 = Inches(1.7), Inches(0.9), Inches(0.3)
    xS = meso_panel.left + Inches(0.35)

    s_data = add_box(slide, xS, yS, w2, h2, "Meso Data"); xS += w2 + g2
    s_core = add_box(slide, xS, yS, w2, h2, "Meso AI Core"); xS += w2 + g2
    s_val  = add_box(slide, xS, yS, w2, h2, "Validated Meso Model"); xS += w2 + g2
    s_oed  = add_diamond(slide, xS, yS - Inches(0.05), Inches(1.2), Inches(1.2), "Meso OED"); xS += Inches(1.2) + g2
    s_exp  = add_box(slide, xS, yS, w2, h2, "Meso Experiments")

    connect_lr(slide, s_data, s_core)
    connect_lr(slide, s_core, s_val)
    connect_lr(slide, s_val, s_oed)
    connect_lr(slide, s_oed, s_exp)
    elbow_feedback(slide, s_exp, s_data, via_y=meso_panel.top + Inches(2.25))

    x1, y1 = bottom_center(m_val)
    x2, y2 = top_center(s_core)
    add_arrow(slide, x1, y1, x2, y2)
    add_label(slide, (x1 + x2) // 2 - Inches(1.2), (y1 + y2) // 2 - Inches(0.4), Inches(2.8), Inches(0.5),
              "Macro data provides target\nfor inverse analysis")

    x1b, y1b = bottom_center(m_val)
    x2b, y2b = top_center(s_val)
    add_arrow(slide, x1b, y1b, x2b, y2b, dashed=True)
    add_label(slide, (x1b + x2b) // 2 - Inches(1.4), (y1b + y2b) // 2, Inches(3.2), Inches(0.5),
              "Predicts macro response /\nProvides homogenized properties as input")

def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide1(prs)
    slide2(prs)

    prs.save("hybrid_models.pptx")
    print("Generated hybrid_models.pptx with 2 slides.")

if __name__ == "__main__":
    build()
