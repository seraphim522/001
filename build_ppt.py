from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
import os

TITLE = "Research Pipeline"
SUBTITLE_CN = "理论 → 模型 → 方法 → 算法 → 代码 → 原型软件 → 软件平台"
STAGES = [
    "Theory", "Models", "Methods", "Algorithms",
    "Code", "Prototype\nSoftware", "Software\nPlatform"
]

# Blue-cyan gradient-like palette (left -> right)
PALETTE = ["#0b63ce", "#1e78d5", "#2e8ddc", "#3ea2e3", "#4eb7ea", "#5fcdf1", "#70e2f8"]

def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

def add_title(slide, prs):
    # Title
    left = Inches(0.6)
    top = Inches(0.35)
    width = prs.slide_width - Inches(1.2)
    height = Inches(1.0)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = TITLE
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 20, 20)
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Subtitle (CN)
    top2 = top + Inches(0.8)
    sub_box = slide.shapes.add_textbox(left, top2, width, Inches(0.7))
    tf2 = sub_box.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.text = SUBTITLE_CN
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(80, 80, 80)
    p2.alignment = PP_ALIGN.CENTER
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_flow(slide, prs):
    margin_lr = Inches(0.6)
    top = Inches(2.1)
    height = Inches(1.6)

    avail_w = prs.slide_width - margin_lr * 2
    n = len(STAGES)
    chevron_w = avail_w / n

    for i, label in enumerate(STAGES):
        left = margin_lr + chevron_w * i
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left, top, chevron_w, height)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = rgb(PALETTE[i % len(PALETTE)])
        line = shape.line
        line.width = Pt(1.25)
        line.color.rgb = RGBColor(255, 255, 255)

        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

def add_notes(slide, prs):
    # Concise one-liners for "说明不同阶段"
    notes = (
        "From hypothesis to deployed platform: each stage validates and de-risks the next.\n"
        "• Theory: rigorous problem + assumptions   • Models: abstractions + metrics\n"
        "• Methods/Algorithms: design + proof/analysis   • Code: reproducible implementation\n"
        "• Prototype: empirical validation   • Platform: productized, scalable"
    )
    left = Inches(0.6)
    top = Inches(4.1)
    width = prs.slide_width - Inches(1.2)
    height = Inches(1.6)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = notes
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(60, 60, 60)
    p.alignment = PP_ALIGN.LEFT
    tf.vertical_anchor = MSO_ANCHOR.TOP

def build():
    prs = Presentation()
    # 16:9 canvas
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # light background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    add_title(slide, prs)
    add_flow(slide, prs)
    add_notes(slide, prs)

    os.makedirs("output", exist_ok=True)
    out = os.path.join("output", "research-pipeline.pptx")
    prs.save(out)
    print(f"Saved -> {out}")

if __name__ == "__main__":
    build()
